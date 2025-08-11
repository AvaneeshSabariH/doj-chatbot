#!/usr/bin/env python3
"""
doj_pipeline.py
Single-file pipeline:
 - DFS crawl (with allowed domains, depth limit)
 - Download HTML/PDF/DOCX/PPTX/XLSX (with cache using ETag/Last-Modified and SHA256)
 - Extract text (html -> text + tables, pdf/docx/pptx/xlsx -> text)
 - Clean & save to ./data/ (text files)
 - Move failed extractions to ./data_failed/
 - Rebuild ChromaDB using Google embeddings (keeps your existing langchain setup)
 - Run a short QA test suite and print context+answer
"""

import os
import sys
import re
import json
import time
import hashlib
import shutil
import argparse
from urllib.parse import urljoin, urlparse
from collections import deque

import requests
from bs4 import BeautifulSoup

# Document extraction libs
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

# Langchain & Chroma (matches your existing tester)
from dotenv import load_dotenv
load_dotenv()
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_community.document_loaders import DirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.prompts import PromptTemplate
from langchain.chains.retrieval import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain

# -----------------------------
# CONFIG
# -----------------------------
START_URLS = ["https://doj.gov.in/"]
ALLOWED_DOMAINS = ["doj.gov.in", "dashboard.doj.gov.in", "scdg.sci.gov.in"]
MAX_DEPTH = 4
DOWNLOAD_FILES = True
CACHE_DIR = "cache"
RAW_DIR = os.path.join(CACHE_DIR, "raw")
DATA_DIR = "data"                # cleaned text output used for ingestion
FAILED_DIR = "data_failed"       # failed extraction files
PERSIST_DIR = "chroma_db"
HASH_CACHE_FILE = os.path.join(CACHE_DIR, "hash_cache.json")
SKIP_UNCHANGED = True
REQUESTS_HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
TEST_QUESTIONS = [
    "What are Fast Track Special Courts?",
    "What is the purpose of the Tele-Law service?",
    "Explain the role of the National Judicial Data Grid.",
    "What is Nyaya Bandhu and who can use it?",
    "What are the objectives of the eCourts project?",
    "What is the Citizens’ Charter?",
    "What does the Gram Nyayalaya scheme aim to do?",
    "List some services offered under eCourt Services.",
    "What is the Memorandum of procedure of appointment of Supreme Court Judges?",
    "What is the function of the National Legal Services Authority (NALSA)?"
]
MIN_TEXT_WORDS = 8   # keep short but meaningful content threshold

# -----------------------------
# UTILITIES
# -----------------------------
def ensure_dirs():
    os.makedirs(CACHE_DIR, exist_ok=True)
    os.makedirs(RAW_DIR, exist_ok=True)
    os.makedirs(DATA_DIR, exist_ok=True)
    os.makedirs(FAILED_DIR, exist_ok=True)

def load_hash_cache():
    if os.path.exists(HASH_CACHE_FILE):
        try:
            with open(HASH_CACHE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_hash_cache(cache):
    with open(HASH_CACHE_FILE, "w", encoding="utf-8") as f:
        json.dump(cache, f, indent=2)

def sha256_bytes(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()

def normalize_url(u: str) -> str:
    # remove fragments, unify trailing slash
    parsed = urlparse(u)
    path = parsed.path or "/"
    scheme = parsed.scheme or "https"
    netloc = parsed.netloc
    new = f"{scheme}://{netloc}{path}"
    return new.rstrip("/")

def allowed_domain(url: str) -> bool:
    parsed = urlparse(url)
    domain = parsed.netloc
    return any(domain.endswith(d) for d in ALLOWED_DOMAINS)

def safe_filename_from_url(url: str) -> str:
    parsed = urlparse(url)
    name = parsed.netloc + parsed.path
    name = re.sub(r"[^a-zA-Z0-9_\-\.]", "_", name)
    return name.strip("_")

# -----------------------------
# CRAWLER + DOWNLOADER
# -----------------------------
def head_request(url):
    try:
        r = requests.head(url, headers=REQUESTS_HEADERS, allow_redirects=True, timeout=10)
        return r
    except Exception:
        return None

def download_if_changed(url, hash_cache):
    """
    Download resource into RAW_DIR if:
      - not present in cache, or
      - ETag/Last-Modified changed, or
      - content SHA changed
    Returns tuple (status, local_path) where status in {"skipped","new","updated","failed"}
    """
    urln = normalize_url(url)
    filename = safe_filename_from_url(urln)
    # keep original extension if possible
    ext = os.path.splitext(urlparse(url).path)[1].lower()
    if not ext:
        ext = ".html"
    local_path = os.path.join(RAW_DIR, filename + ext)

    # HEAD for ETag / Last-Modified
    try:
        head = head_request(url)
    except Exception:
        head = None

    etag = head.headers.get("ETag") if head is not None else None
    last_mod = head.headers.get("Last-Modified") if head is not None else None

    cache_entry = hash_cache.get(urln, {})
    cached_etag = cache_entry.get("etag")
    cached_last = cache_entry.get("last_modified")
    cached_hash = cache_entry.get("sha256")

    # If file exists and skip unchanged
    if SKIP_UNCHANGED and os.path.exists(local_path):
        # Quick check via headers
        if etag and cached_etag and etag == cached_etag:
            print(f"[Cache] Skipping unchanged: {local_path}")
            return "skipped", local_path
        if last_mod and cached_last and last_mod == cached_last:
            print(f"[Cache] Skipping unchanged (last-mod): {local_path}")
            return "skipped", local_path
        # fallback: compare content hash
        try:
            with open(local_path, "rb") as f:
                existing = f.read()
            if cached_hash and sha256_bytes(existing) == cached_hash:
                print(f"[Cache] Skipping unchanged (hash): {local_path}")
                return "skipped", local_path
        except Exception:
            pass

    # Do GET
    try:
        r = requests.get(url, headers=REQUESTS_HEADERS, allow_redirects=True, timeout=20)
        r.raise_for_status()
        content = r.content
        # write to local_path
        with open(local_path, "wb") as f:
            f.write(content)
        new_hash = sha256_bytes(content)
        new_entry = {"etag": etag, "last_modified": last_mod, "sha256": new_hash, "path": local_path}
        prev_hash = cache_entry.get("sha256")
        if prev_hash and prev_hash != new_hash:
            print(f"[Updated] Downloaded: {local_path}")
            hash_cache[urln] = new_entry
            return "updated", local_path
        else:
            print(f"[New] Downloaded: {local_path}")
            hash_cache[urln] = new_entry
            return "new", local_path
    except Exception as e:
        print(f"[Fail] Download failed: {url} -- {e}")
        return "failed", None

def crawl_and_download(seeds, max_depth, hash_cache):
    visited = set()
    q = deque()
    for s in seeds:
        q.append((normalize_url(s), 0))
    downloaded = []
    while q:
        url, depth = q.pop()
        if depth > max_depth:
            continue
        if url in visited:
            continue
        visited.add(url)
        if not allowed_domain(url):
            continue
        # skip anchors
        if url.endswith("#SkipContent") or "#" in url:
            url = url.split("#")[0]
        print(f"Crawling: {url}")
        status, local = download_if_changed(url, hash_cache) if DOWNLOAD_FILES else ("skipped", None)
        if status in ("new", "updated"):
            downloaded.append(local)
        # if HTML add children links
        try:
            r = requests.get(url, headers=REQUESTS_HEADERS, timeout=10)
            r.raise_for_status()
            if "text/html" in r.headers.get("Content-Type", ""):
                soup = BeautifulSoup(r.text, "html.parser")
                for a in soup.find_all("a", href=True):
                    href = a["href"]
                    # build absolute
                    href_abs = urljoin(url, href)
                    href_abs = href_abs.split("#")[0]  # strip fragment
                    # normalize and enqueue
                    norm = normalize_url(href_abs)
                    if norm not in visited and allowed_domain(norm):
                        q.append((norm, depth + 1))
        except Exception:
            # ignore parse errors for link extraction
            pass
    return downloaded

# -----------------------------
# EXTRACTION / CLEANING
# -----------------------------
def ensure_text_minimum(text: str) -> bool:
    words = re.findall(r"\w+", text)
    return len(words) >= MIN_TEXT_WORDS

def clean_extracted_text(text: str) -> str:
    # basic normalization: collapse whitespace, remove repeated header/footer patterns
    text = re.sub(r"\r\n|\r", "\n", text)
    # remove long runs of whitespace and control chars
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # strip common website boilerplate lines aggressively
    boilerplate_patterns = [
        r"Content\s+Owned\s+by",
        r"Developed\s+and\s+hosted\s+by",
        r"Last\s+Updated",
        r"Disclaimer",
        r"All\s+rights\s+reserved",
        r"Visitor\s+Summary"
    ]
    lines = []
    for line in text.splitlines():
        if any(re.search(pat, line, re.IGNORECASE) for pat in boilerplate_patterns):
            continue
        lines.append(line.strip())
    cleaned = "\n".join([ln for ln in lines if ln])
    # final whitespace normalization
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned

def extract_from_pdf(path):
    try:
        txt = pdf_extract_text(path)
        return txt
    except Exception as e:
        raise RuntimeError(f"pdf extract failed: {e}")

def extract_from_docx(path):
    try:
        doc = DocxDocument(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text)
        return text
    except Exception as e:
        raise RuntimeError(f"docx extract failed: {e}")

def extract_from_pptx(path):
    try:
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out)
    except Exception as e:
        raise RuntimeError(f"pptx extract failed: {e}")

def extract_from_xlsx(path):
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        rows = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows.append(f"Sheet: {name}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join("" if v is None else str(v) for v in row)
                rows.append(row_text)
        return "\n".join(rows)
    except Exception as e:
        raise RuntimeError(f"xlsx extract failed: {e}")

def extract_from_html(path_or_url, is_local=True):
    # path_or_url can be file path or remote url; prefer to parse file if local
    try:
        if is_local:
            with open(path_or_url, "r", encoding="utf-8", errors="ignore") as f:
                html = f.read()
        else:
            r = requests.get(path_or_url, headers=REQUESTS_HEADERS, timeout=10)
            r.raise_for_status()
            html = r.text
        soup = BeautifulSoup(html, "html.parser")
        for s in soup(["script", "style", "noscript", "nav", "footer", "header"]):
            s.decompose()
        text_parts = []
        # collect headings, paragraphs, list items
        for tag in soup.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
            t = tag.get_text(separator=" ", strip=True)
            if t:
                text_parts.append(t)
        # collect tables separately and append as TSV
        table_parts = []
        for table in soup.find_all("table"):
            rows = []
            for tr in table.find_all("tr"):
                cols = [td.get_text(separator=" ", strip=True) for td in tr.find_all(["th", "td"])]
                rows.append("\t".join(cols))
            if rows:
                table_parts.append("\n".join(rows))
        text = "\n\n".join(text_parts)
        if table_parts:
            text += "\n\nTABLES:\n" + "\n\n".join(table_parts)
        return text
    except Exception as e:
        raise RuntimeError(f"html extract failed: {e}")

def process_downloaded_files(downloaded_paths):
    """
    For each local downloaded file, detect type by extension and try to extract text.
    Write cleaned text files to DATA_DIR. Move failures to FAILED_DIR.
    """
    kept = []
    failed = []
    for path in downloaded_paths:
        if not path or not os.path.exists(path):
            continue
        try:
            lower = path.lower()
            if lower.endswith(".pdf"):
                raw = extract_from_pdf(path)
            elif lower.endswith(".docx"):
                raw = extract_from_docx(path)
            elif lower.endswith(".pptx"):
                raw = extract_from_pptx(path)
            elif lower.endswith(".xlsx") or lower.endswith(".xls"):
                raw = extract_from_xlsx(path)
            elif lower.endswith(".html") or lower.endswith(".htm") or lower.endswith(".php") or lower.endswith(".aspx") or lower.endswith(".txt"):
                # many pages are saved as html; if extension is .txt and content is HTML, still attempt
                raw = extract_from_html(path, is_local=True)
            else:
                # fallback: if it's a binary but we saved with .html ext, attempt html parse
                raw = extract_from_html(path, is_local=True)
            cleaned = clean_extracted_text(raw)
            if not ensure_text_minimum(cleaned):
                # if very short (but might still contain useful single-line definitions), keep if meaningful
                # decide heuristics: keep if contains keywords
                keywords = ["Tele-Law", "Tele Law", "eFiling", "Fast Track", "NJDG", "eCourts", "NALSA"]
                if any(k.lower() in cleaned.lower() for k in keywords) and len(cleaned) > 10:
                    pass
                else:
                    raise RuntimeError("content too short or boilerplate")
            # write to DATA_DIR, filename derived from original
            fname = os.path.basename(path)
            outname = os.path.splitext(fname)[0] + ".txt"
            outpath = os.path.join(DATA_DIR, outname)
            with open(outpath, "w", encoding="utf-8") as f:
                f.write(cleaned)
            kept.append(outpath)
        except Exception as e:
            print(f"[Fail] Extract failed for {path}: {e}")
            # move to FAILED_DIR
            try:
                dest = os.path.join(FAILED_DIR, os.path.basename(path))
                shutil.move(path, dest)
            except Exception:
                pass
            failed.append(path)
    return kept, failed

# -----------------------------
# EMBEDDING & RAG (keeps your previous setup)
# -----------------------------
def build_vectorstore_and_chain():
    # --- 1. LOAD DOCUMENTS ---
    loader = DirectoryLoader(DATA_DIR, glob="**/*.txt")
    docs = loader.load()
    print(f"✅ Loaded {len(docs)} document(s).")

    # --- 2. CHUNK DOCUMENTS ---
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = text_splitter.split_documents(docs)
    print(f"✅ Documents split into {len(splits)} chunks.")

    # --- 3. CREATE AND POPULATE VECTOR STORE (CHROMA DB) ---
    print("⏳ Creating vector store from documents...")
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    # wipe existing persistence dir to ensure fresh rebuild
    if os.path.exists(PERSIST_DIR):
        try:
            shutil.rmtree(PERSIST_DIR)
        except Exception:
            pass
    vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings, persist_directory=PERSIST_DIR)
    print("✅ Vector store created successfully.")

    # --- 4. CREATE THE RETRIEVAL CHAIN ---
    llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.3)

    prompt = PromptTemplate(
        template="""You are an assistant for the Department of Justice, India.
Answer the user's question based only on the provided context.
If the information is not in the context, say so.

Context:
{context}

Question:
{input}
""",
        input_variables=["context", "input"],
    )

    document_chain = create_stuff_documents_chain(llm, prompt)
    retriever = vectorstore.as_retriever()
    retrieval_chain = create_retrieval_chain(retriever, document_chain)
    print("✅ Retrieval chain created.")
    return retrieval_chain

# -----------------------------
# QA TEST
# -----------------------------
def run_tests(chain):
    print("\n--- DOJ RAG QA Test Suite ---")
    for q in TEST_QUESTIONS:
        print(f"\nQuestion: {q}")
        try:
            res = chain.invoke({"input": q})
            ans = res.get("answer") or res.get("output_text") or str(res)
            print("Answer:", ans)
        except Exception as e:
            print("⚠️ Error while answering:", e)

# -----------------------------
# MAIN
# -----------------------------
def main():
    ensure_dirs()
    hash_cache = load_hash_cache()

    print("Starting crawl & download...")
    downloaded = crawl_and_download(START_URLS, MAX_DEPTH, hash_cache)
    save_hash_cache(hash_cache)

    print(f"\nProcessing {len(downloaded)} downloaded files...")
    kept, failed = process_downloaded_files(downloaded)
    print(f"✅ Extracted and kept: {len(kept)} files. Failed: {len(failed)} files (moved to {FAILED_DIR}).")

    print("\nBuilding vector store and retrieval chain...")
    chain = build_vectorstore_and_chain()

    print("\nRunning QA tests...")
    run_tests(chain)

    print("\nPipeline finished.")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="DoJ pipeline: crawl/extract/embed/test")
    parser.add_argument("--skip-download", action="store_true", help="Skip crawling/downloading; only run extraction/embedding from cached files")
    args = parser.parse_args()

    if args.skip_download:
        # if skipping download, gather all files already in RAW_DIR
        ensure_dirs()
        hash_cache = load_hash_cache()
        downloaded = [os.path.join(RAW_DIR, f) for f in os.listdir(RAW_DIR)]
        kept, failed = process_downloaded_files(downloaded)
        chain = build_vectorstore_and_chain()
        run_tests(chain)
    else:
        main()
