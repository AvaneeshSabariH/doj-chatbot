import os
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import re

BASE_URL = "https://doj.gov.in/"
VISITED = set()
OUTPUT_DIR = "data"

os.makedirs(OUTPUT_DIR, exist_ok=True)

def clean_text(text):
    text = re.sub(r"\s+", " ", text)
    return text.strip()

def save_text(content, filename, subfolder="misc"):
    folder_path = os.path.join(OUTPUT_DIR, subfolder)
    os.makedirs(folder_path, exist_ok=True)
    filepath = os.path.join(folder_path, filename + ".txt")
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)

def process_pdf(url, filename):
    r = requests.get(url)
    path = os.path.join(OUTPUT_DIR, "pdf", filename + ".pdf")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(r.content)
    try:
        text = pdf_extract_text(path)
        save_text(text, filename, "pdf")
    except Exception as e:
        print(f"⚠️ PDF parse failed for {filename}: {e}")

def process_docx(url, filename):
    r = requests.get(url)
    path = os.path.join(OUTPUT_DIR, "docx", filename + ".docx")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(r.content)
    try:
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        save_text(text, filename, "docx")
    except Exception as e:
        print(f"⚠️ DOCX parse failed for {filename}: {e}")

def process_pptx(url, filename):
    r = requests.get(url)
    path = os.path.join(OUTPUT_DIR, "pptx", filename + ".pptx")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(r.content)
    try:
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        save_text("\n".join(text_runs), filename, "pptx")
    except Exception as e:
        print(f"⚠️ PPTX parse failed for {filename}: {e}")

def process_xlsx(url, filename):
    r = requests.get(url)
    path = os.path.join(OUTPUT_DIR, "xlsx", filename + ".xlsx")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(r.content)
    try:
        wb = load_workbook(path)
        all_text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_data = "\t".join(str(cell) if cell else "" for cell in row)
                all_text.append(row_data)
        save_text("\n".join(all_text), filename, "xlsx")
    except Exception as e:
        print(f"⚠️ XLSX parse failed for {filename}: {e}")

def process_html(url, filename):
    r = requests.get(url)
    soup = BeautifulSoup(r.text, "html.parser")
    
    # Extract visible text
    for script in soup(["script", "style"]):
        script.decompose()
    
    text = clean_text(soup.get_text(separator="\n"))
    
    # Extract tables
    tables_text = []
    for table in soup.find_all("table"):
        for row in table.find_all("tr"):
            cols = [clean_text(col.get_text()) for col in row.find_all(["td", "th"])]
            tables_text.append("\t".join(cols))
    if tables_text:
        text += "\n\nTABLE DATA:\n" + "\n".join(tables_text)

    if len(text) > 200:
        save_text(text, filename, "html")

def crawl(url):
    if url in VISITED:
        return
    VISITED.add(url)
    print(f"Crawling: {url}")

    try:
        r = requests.get(url, timeout=10)
    except Exception as e:
        print(f"⚠️ Request failed: {e}")
        return
    
    content_type = r.headers.get("Content-Type", "")
    parsed = urlparse(url)
    filename = re.sub(r"[^a-zA-Z0-9_-]", "_", parsed.path.strip("/")) or "index"

    if "application/pdf" in content_type:
        process_pdf(url, filename)
    elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type:
        process_docx(url, filename)
    elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in content_type:
        process_pptx(url, filename)
    elif "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" in content_type:
        process_xlsx(url, filename)
    elif "text/html" in content_type:
        process_html(url, filename)
        soup = BeautifulSoup(r.text, "html.parser")
        for a in soup.find_all("a", href=True):
            link = urljoin(url, a["href"])
            if link.startswith(BASE_URL):
                crawl(link)

if __name__ == "__main__":
    crawl(BASE_URL)
